#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
XONIMET 2026 - Extractor Universal de Metadatos
Modo: Interfaz Gráfica (Flask) o Terminal (CLI)
Optimizado para equipos de bajos recursos
Desarrollador: Darian Alberto Camacho Salas
Organización: XONIDU
"""

import os
import sys
import json
import datetime
import hashlib
import warnings
from pathlib import Path
warnings.filterwarnings('ignore')

# ============================================================================
# Verificación de dependencias con mensajes claros
# ============================================================================
def check_import(module_name, package_name, error_msg):
    try:
        __import__(module_name)
        return True
    except ImportError:
        print(f"[ERROR] {error_msg}")
        print(f"        Ejecuta: pip install {package_name}")
        return False

# Dependencias necesarias
DEPENDENCIAS = [
    ('PIL', 'pillow', 'Pillow no instalado (necesario para imágenes)'),
    ('mutagen', 'mutagen', 'Mutagen no instalado (necesario para audio)'),
    ('PyPDF2', 'pypdf2', 'PyPDF2 no instalado (necesario para PDF)'),
    ('docx', 'python-docx', 'python-docx no instalado (necesario para Word)'),
    ('openpyxl', 'openpyxl', 'openpyxl no instalado (necesario para Excel)'),
    ('pptx', 'python-pptx', 'python-pptx no instalado (necesario para PowerPoint)'),
    ('exifread', 'exifread', 'exifread no instalado (necesario para EXIF)'),
    ('reportlab', 'reportlab', 'reportlab no instalado (necesario para PDF)')
]

# Dependencias opcionales (solo para modo gráfico)
DEPENDENCIAS_GRAFICAS = [
    ('flask', 'flask', 'Flask no instalado (necesario para interfaz gráfica)'),
    ('ffmpeg', 'ffmpeg-python', 'ffmpeg-python no instalado (necesario para videos)')
]

def verificar_dependencias(modo_grafico=False):
    """Verifica dependencias según el modo seleccionado"""
    ok = True
    
    # Dependencias básicas
    for module, package, msg in DEPENDENCIAS:
        if not check_import(module, package, msg):
            ok = False
    
    # Dependencias para modo gráfico
    if modo_grafico:
        for module, package, msg in DEPENDENCIAS_GRAFICAS:
            if not check_import(module, package, msg):
                ok = False
    
    return ok

# ============================================================================
# Colores para terminal
# ============================================================================
class Colors:
    GREEN = '\033[92m'
    YELLOW = '\033[93m'
    RED = '\033[91m'
    BLUE = '\033[94m'
    PURPLE = '\033[95m'
    CYAN = '\033[96m'
    END = '\033[0m'
    BOLD = '\033[1m'

def print_banner():
    banner = f"""
{Colors.PURPLE}{Colors.BOLD}═══════════════════════════════════════════════════════════
                    XONIMET 2026 v3.0                    
              Extractor Universal de Metadatos            
              Extrae informacion de: Fotos, Audio,        
              Video, Documentos y mas                     
                                                          
              [1] Modo Terminal (CLI)                     
              [2] Modo Grafico (Web)                      
                                                          
              Desarrollado por: Darian Alberto            
              Camacho Salas                               
              #Somos XONIDU
═══════════════════════════════════════════════════════════{Colors.END}
    """
    print(banner)

def mostrar_ayuda():
    ayuda = f"""
{Colors.BOLD}USO DE XONIMET:{Colors.END}

  python xonimet.py [opcion]

{Colors.BOLD}OPCIONES:{Colors.END}

  [1]  Modo Terminal (CLI) - Interactivo por consola
  [2]  Modo Grafico (Web)  - Interfaz web en http://localhost:5000

{Colors.BOLD}EJEMPLOS:{Colors.END}

  Modo Terminal:
    python xonimet.py 1

  Modo Grafico:
    python xonimet.py 2

{Colors.BOLD}MODO TERMINAL (CLI):{Colors.END}

  Una vez iniciado, tendras un menu con las opciones:
    1. Seleccionar archivo
    2. Analizar archivo actual
    3. Guardar resultados en JSON
    4. Generar reporte PDF
    5. Ver informacion del archivo
    6. Cambiar archivo
    7. Ayuda
    0. Salir

{Colors.BOLD}MODO GRAFICO (WEB):{Colors.END}

  Abre tu navegador y ve a: http://localhost:5000
  Sube archivos y visualiza metadatos en una interfaz amigable.
    """
    print(ayuda)

# ============================================================================
# Clase Xonimet (núcleo del extractor)
# ============================================================================
class Xonimet:
    def __init__(self, file_path=None):
        self.file_path = Path(file_path) if file_path else None
        self.metadata = {}
        self.file_content = None
    
    def set_file(self, file_path):
        self.file_path = Path(file_path)
        self.metadata = {}
        self.file_content = None
    
    def set_file_content(self, file_content, filename):
        self.file_content = file_content
        self.file_path = Path(filename)
        self.metadata = {}
    
    def get_config_dir(self):
        return os.path.join(os.path.expanduser("~"), '.xonimet')
    
    def ensure_config_dir(self):
        config_dir = self.get_config_dir()
        os.makedirs(config_dir, exist_ok=True)
        return config_dir
    
    def extract_all(self):
        if not self.file_path:
            return {'error': 'No se ha seleccionado ningún archivo'}
        
        self.metadata = {
            'archivo': {
                'nombre': self.file_path.name,
                'tamaño_bytes': len(self.file_content) if self.file_content else self.file_path.stat().st_size,
                'tamaño_formateado': self._format_bytes(len(self.file_content) if self.file_content else self.file_path.stat().st_size),
                'extension': self.file_path.suffix.lower(),
                'tipo_mime': self._get_mime_type(),
                'hashes': self._calculate_hashes()
            },
            'metadatos_especificos': self._extract_specific_metadata()
        }
        return self.metadata
    
    def _format_bytes(self, bytes):
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if bytes < 1024.0:
                return f"{bytes:.2f} {unit}"
            bytes /= 1024.0
        return f"{bytes:.2f} PB"
    
    def _get_mime_type(self):
        import mimetypes
        mime_type, _ = mimetypes.guess_type(str(self.file_path))
        return mime_type or 'desconocido'
    
    def _calculate_hashes(self):
        hashes = {}
        try:
            data = self.file_content if self.file_content else open(self.file_path, 'rb').read()
            hashes['md5'] = hashlib.md5(data).hexdigest()
            hashes['sha1'] = hashlib.sha1(data).hexdigest()
            hashes['sha256'] = hashlib.sha256(data).hexdigest()
        except:
            hashes['error'] = 'No se pudo calcular hashes'
        return hashes
    
    def _extract_image_metadata(self):
        metadata = {}
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS
            import exifread
            from io import BytesIO
            
            img_data = self.file_content if self.file_content else open(self.file_path, 'rb').read()
            img = Image.open(BytesIO(img_data))
            metadata.update({
                'dimensiones': f"{img.width} x {img.height}",
                'modo': img.mode,
                'formato': img.format,
                'info_basica': {
                    'ancho': img.width,
                    'alto': img.height,
                    'proporcion': round(img.width / img.height, 2)
                }
            })
            
            if hasattr(img, '_getexif') and img._getexif():
                exif = {}
                for tag_id, value in img._getexif().items():
                    tag = TAGS.get(tag_id, tag_id)
                    if isinstance(value, bytes):
                        try:
                            value = value.decode('utf-8', errors='ignore')
                        except:
                            value = str(value)
                    exif[tag] = str(value)
                metadata['exif'] = exif
            
            with BytesIO(img_data) as f:
                tags = exifread.process_file(f, details=True)
                if tags:
                    metadata['exif_detallado'] = {str(k): str(v) for k, v in tags.items()}
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_audio_metadata(self):
        metadata = {}
        try:
            import mutagen
            audio = mutagen.File(self.file_path)
            if audio:
                metadata['formato'] = type(audio).__name__
                metadata['duracion_segundos'] = audio.info.length
                metadata['duracion_formateado'] = str(datetime.timedelta(seconds=int(audio.info.length)))
                if hasattr(audio.info, 'bitrate'):
                    metadata['bitrate'] = f"{audio.info.bitrate // 1000} kbps"
                if hasattr(audio.info, 'sample_rate'):
                    metadata['frecuencia_muestreo'] = f"{audio.info.sample_rate} Hz"
                if hasattr(audio, 'tags') and audio.tags:
                    tags = {}
                    for key, value in audio.tags.items():
                        if value:
                            tags[key] = str(value[0]) if isinstance(value, list) else str(value)
                    metadata['etiquetas'] = tags
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_video_metadata(self):
        metadata = {}
        try:
            import ffmpeg
            
            temp_path = f"/tmp/{self.file_path.name}"
            with open(temp_path, 'wb') as f:
                f.write(self.file_content if self.file_content else open(self.file_path, 'rb').read())
            
            probe = ffmpeg.probe(temp_path)
            os.remove(temp_path)
            
            if 'format' in probe:
                fmt = probe['format']
                metadata['formato'] = {
                    'nombre': fmt.get('format_name'),
                    'duracion': fmt.get('duration'),
                    'bitrate': fmt.get('bit_rate'),
                    'tamaño': fmt.get('size'),
                    'tags': fmt.get('tags', {})
                }
            
            streams = []
            for stream in probe.get('streams', []):
                stream_info = {'tipo': stream.get('codec_type'), 'codec': stream.get('codec_name')}
                if stream['codec_type'] == 'video':
                    stream_info.update({
                        'resolucion': f"{stream.get('width')}x{stream.get('height')}",
                        'fps': eval(stream.get('r_frame_rate', '0/1')),
                        'pixeles': stream.get('pix_fmt')
                    })
                elif stream['codec_type'] == 'audio':
                    stream_info.update({
                        'canales': stream.get('channels'),
                        'frecuencia': stream.get('sample_rate')
                    })
                streams.append(stream_info)
            metadata['streams'] = streams
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_pdf_metadata(self):
        metadata = {}
        try:
            from PyPDF2 import PdfReader
            pdf = PdfReader(self.file_path)
            metadata.update({
                'paginas': len(pdf.pages),
                'encriptado': pdf.is_encrypted,
                'metadatos': dict(pdf.metadata) if pdf.metadata else {}
            })
            if len(pdf.pages) > 0:
                first_page = pdf.pages[0]
                text = first_page.extract_text()
                metadata['primeras_palabras'] = text[:200] + '...' if len(text) > 200 else text
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_docx_metadata(self):
        metadata = {}
        try:
            import docx
            doc = docx.Document(self.file_path)
            core_props = doc.core_properties
            metadata.update({
                'autor': core_props.author,
                'creador': core_props.created,
                'modificado_por': core_props.last_modified_by,
                'fecha_modificacion': core_props.modified,
                'titulo': core_props.title,
                'asunto': core_props.subject,
                'palabras_clave': core_props.keywords,
                'categoria': core_props.category,
                'comentarios': core_props.comments,
                'parrafos': len(doc.paragraphs),
                'tablas': len(doc.tables)
            })
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_xlsx_metadata(self):
        metadata = {}
        try:
            from openpyxl import load_workbook
            wb = load_workbook(self.file_path, data_only=True)
            metadata.update({
                'hojas': wb.sheetnames,
                'hojas_activas': len(wb.sheetnames),
                'propiedades': {
                    'creador': wb.properties.creator,
                    'creado': str(wb.properties.created) if wb.properties.created else None,
                    'modificado': str(wb.properties.modified) if wb.properties.modified else None,
                    'titulo': wb.properties.title
                }
            })
            sheets_info = {}
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheets_info[sheet_name] = {
                    'filas': sheet.max_row,
                    'columnas': sheet.max_column,
                    'celdas_con_datos': sum(1 for row in sheet.iter_rows() for cell in row if cell.value)
                }
            metadata['detalle_hojas'] = sheets_info
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_pptx_metadata(self):
        metadata = {}
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            core_props = prs.core_properties
            metadata.update({
                'diapositivas': len(prs.slides),
                'autor': core_props.author,
                'creado': str(core_props.created) if core_props.created else None,
                'modificado': str(core_props.modified) if core_props.modified else None,
                'titulo': core_props.title,
                'asunto': core_props.subject
            })
            slide_stats = {'texto': 0, 'imagenes': 0, 'tablas': 0, 'graficos': 0}
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_stats['texto'] += 1
                    if hasattr(shape, 'image'):
                        slide_stats['imagenes'] += 1
                    if shape.has_table:
                        slide_stats['tablas'] += 1
                    if hasattr(shape, 'chart'):
                        slide_stats['graficos'] += 1
            metadata['estadisticas_diapositivas'] = slide_stats
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_text_metadata(self):
        metadata = {}
        try:
            content = self.file_content if self.file_content else open(self.file_path, 'r', encoding='utf-8', errors='ignore').read()
            lines = content.split('\n')
            metadata.update({
                'lineas': len(lines),
                'palabras': len(content.split()),
                'caracteres': len(content),
                'caracteres_sin_espacios': len(content.replace(' ', '').replace('\n', '').replace('\t', '')),
                'primeras_10_lineas': lines[:10] if len(lines) > 10 else lines
            })
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_specific_metadata(self):
        ext = self.file_path.suffix.lower()
        if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.heic']:
            return self._extract_image_metadata()
        elif ext in ['.mp3', '.flac', '.wav', '.ogg', '.m4a', '.aac', '.wma']:
            return self._extract_audio_metadata()
        elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v']:
            return self._extract_video_metadata()
        elif ext == '.pdf':
            return self._extract_pdf_metadata()
        elif ext in ['.docx', '.doc']:
            return self._extract_docx_metadata()
        elif ext in ['.xlsx', '.xls']:
            return self._extract_xlsx_metadata()
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx_metadata()
        elif ext in ['.txt', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.md']:
            return self._extract_text_metadata()
        else:
            return {'mensaje': 'Tipo de archivo no soportado'}
    
    def print_metadata(self, metadata=None):
        if metadata is None:
            metadata = self.metadata
        
        if not metadata:
            print(f"{Colors.RED}No hay metadatos para mostrar{Colors.END}")
            return
        
        print(f"\n{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}")
        print(f"{Colors.CYAN}{Colors.BOLD}METADATOS DEL ARCHIVO{Colors.END}")
        print(f"{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}")
        
        print(f"\n{Colors.GREEN}{Colors.BOLD}INFORMACION BASICA:{Colors.END}")
        archivo = metadata.get('archivo', {})
        for key, value in archivo.items():
            if key != 'hashes':
                print(f"  {Colors.YELLOW}•{Colors.END} {key.replace('_', ' ').title()}: {value}")
        
        if 'hashes' in archivo:
            print(f"\n{Colors.GREEN}{Colors.BOLD}HASHES:{Colors.END}")
            for algo, hash_value in archivo['hashes'].items():
                print(f"  {Colors.YELLOW}•{Colors.END} {algo.upper()}: {hash_value}")
        
        if 'metadatos_especificos' in metadata and metadata['metadatos_especificos']:
            print(f"\n{Colors.GREEN}{Colors.BOLD}METADATOS ESPECIFICOS:{Colors.END}")
            spec = metadata['metadatos_especificos']
            
            if 'error' in spec:
                print(f"  {Colors.RED}⚠ {spec['error']}{Colors.END}")
            else:
                for key, value in spec.items():
                    if isinstance(value, dict):
                        print(f"\n  {Colors.CYAN}• {key.replace('_', ' ').title()}:{Colors.END}")
                        for subkey, subvalue in value.items():
                            if subvalue:
                                print(f"    {Colors.YELLOW}•{Colors.END} {subkey}: {subvalue}")
                    else:
                        print(f"  {Colors.YELLOW}•{Colors.END} {key.replace('_', ' ').title()}: {value}")
        
        print(f"\n{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}")
    
    def generate_pdf_report(self, output_path=None):
        if not self.metadata:
            return None
        
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib import colors
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.enums import TA_CENTER
            
            self.ensure_config_dir()
            
            if output_path is None:
                base_name = Path(self.metadata['archivo']['nombre']).stem
                config_dir = self.get_config_dir()
                if os.access(config_dir, os.W_OK):
                    output_path = os.path.join(config_dir, f"{base_name}_reporte.pdf")
                else:
                    output_path = f"{base_name}_reporte.pdf"
            
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontSize=20, textColor=colors.HexColor('#1a5490'), alignment=TA_CENTER, spaceAfter=30)
            heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=14, textColor=colors.HexColor('#2c7ab1'), spaceAfter=12, spaceBefore=20)
            normal_style = ParagraphStyle('CustomNormal', parent=styles['Normal'], fontSize=10, spaceAfter=6)
            
            story.append(Paragraph("Reporte de Metadatos", title_style))
            
            archivo = self.metadata.get('archivo', {})
            story.append(Paragraph(f"Archivo: {archivo.get('nombre', 'Desconocido')}", heading_style))
            
            archivo_data = []
            for key, value in archivo.items():
                if key != 'hashes':
                    archivo_data.append([Paragraph(f"<b>{key.replace('_', ' ').title()}</b>", normal_style), Paragraph(str(value), normal_style)])
            
            if archivo_data:
                table = Table(archivo_data, colWidths=[2*inch, 3.5*inch])
                table.setStyle(TableStyle([('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#e8f0f8')), ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#c0c0c0')), ('PADDING', (0, 0), (-1, -1), 6)]))
                story.append(table)
            
            if 'hashes' in archivo:
                story.append(Spacer(1, 0.2*inch))
                story.append(Paragraph("Hashes de Seguridad", heading_style))
                hash_data = []
                for algo, hash_value in archivo['hashes'].items():
                    hash_data.append([Paragraph(f"<b>{algo.upper()}</b>", normal_style), Paragraph(str(hash_value), normal_style)])
                if hash_data:
                    hash_table = Table(hash_data, colWidths=[1.2*inch, 4.3*inch])
                    hash_table.setStyle(TableStyle([('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#e8f0f8')), ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#c0c0c0')), ('PADDING', (0, 0), (-1, -1), 6)]))
                    story.append(hash_table)
            
            spec = self.metadata.get('metadatos_especificos', {})
            if spec and 'error' not in spec:
                story.append(Spacer(1, 0.3*inch))
                story.append(Paragraph("Metadatos Específicos", heading_style))
                for key, value in spec.items():
                    if isinstance(value, dict):
                        story.append(Paragraph(f"<b>{key.replace('_', ' ').title()}</b>", styles['Heading3']))
                        sub_data = []
                        for subkey, subvalue in value.items():
                            if subvalue:
                                sub_data.append([Paragraph(f"<i>{subkey.replace('_', ' ').title()}</i>", normal_style), Paragraph(str(subvalue), normal_style)])
                        if sub_data:
                            sub_table = Table(sub_data, colWidths=[2*inch, 3.5*inch])
                            sub_table.setStyle(TableStyle([('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d0d0d0')), ('PADDING', (0, 0), (-1, -1), 4)]))
                            story.append(sub_table)
                    else:
                        if value:
                            story.append(Paragraph(f"{key.replace('_', ' ').title()}: {value}", normal_style))
            
            story.append(Spacer(1, 0.5*inch))
            footer_text = f"Reporte generado por XONIMET 2026 - {datetime.datetime.now().strftime('%d/%m/%Y %H:%M:%S')}"
            story.append(Paragraph(footer_text, ParagraphStyle('Footer', parent=styles['Normal'], fontSize=8, textColor=colors.HexColor('#888888'), alignment=TA_CENTER)))
            
            doc.build(story)
            return output_path
        except Exception as e:
            print(f"{Colors.RED}Error generando PDF: {e}{Colors.END}")
            return None

# ============================================================================
# Modo Terminal (CLI)
# ============================================================================
def clear_screen():
    os.system('cls' if os.name == 'nt' else 'clear')

def print_menu():
    menu = f"""
{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════
                   XONIMET 2026 - MODO TERMINAL          
              Extractor Universal de Metadatos           
═══════════════════════════════════════════════════════════{Colors.END}

{Colors.BOLD}MENU PRINCIPAL:{Colors.END}
  {Colors.CYAN}[1]{Colors.END} Seleccionar archivo para analizar
  {Colors.CYAN}[2]{Colors.END} Analizar archivo actual
  {Colors.CYAN}[3]{Colors.END} Guardar resultados en JSON
  {Colors.CYAN}[4]{Colors.END} Generar reporte PDF
  {Colors.CYAN}[5]{Colors.END} Ver informacion del archivo actual
  {Colors.CYAN}[6]{Colors.END} Cambiar archivo
  {Colors.CYAN}[7]{Colors.END} Ayuda / Formatos soportados
  {Colors.CYAN}[8]{Colors.END} Limpiar pantalla
  {Colors.CYAN}[0]{Colors.END} Salir

{Colors.YELLOW}═══════════════════════════════════════════════════════════{Colors.END}
"""
    print(menu)

def select_file():
    print(f"\n{Colors.CYAN}SELECCIONAR ARCHIVO{Colors.END}")
    print(f"{Colors.YELLOW}Escribe la ruta del archivo (o 'cancel' para volver):{Colors.END}")
    
    while True:
        file_path = input(f"{Colors.GREEN}→{Colors.END} ").strip()
        if file_path.lower() == 'cancel':
            return None
        if not file_path:
            continue
        file_path = os.path.expanduser(file_path)
        if os.path.exists(file_path):
            return file_path
        else:
            print(f"{Colors.RED}El archivo no existe. Intenta de nuevo:{Colors.END}")

def print_help():
    help_text = f"""
{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════
AYUDA - FORMATOS SOPORTADOS
═══════════════════════════════════════════════════════════{Colors.END}

{Colors.GREEN}IMAGENES:{Colors.END}
  • .jpg, .jpeg, .png, .gif, .bmp, .tiff, .webp, .heic
  {Colors.YELLOW}→{Colors.END} EXIF, dimensiones, GPS, modelo camara, fecha

{Colors.GREEN}AUDIO:{Colors.END}
  • .mp3, .flac, .wav, .ogg, .m4a, .aac, .wma
  {Colors.YELLOW}→{Colors.END} Duracion, bitrate, etiquetas ID3, artista, album

{Colors.GREEN}VIDEO:{Colors.END}
  • .mp4, .avi, .mov, .mkv, .wmv, .flv, .webm
  {Colors.YELLOW}→{Colors.END} Resolucion, codecs, fps, streams

{Colors.GREEN}DOCUMENTOS:{Colors.END}
  • .pdf: paginas, autor, titulo, metadatos
  • .docx, .doc: autor, fechas, estadisticas
  • .xlsx, .xls: hojas, celdas, propiedades
  • .pptx, .ppt: diapositivas, estadisticas

{Colors.GREEN}TEXTO:{Colors.END}
  • .txt, .csv, .json, .xml, .html, .css, .js, .py, .md
  {Colors.YELLOW}→{Colors.END} Lineas, palabras, caracteres

{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}
"""
    print(help_text)

def save_to_json(metadata):
    if not metadata:
        print(f"{Colors.RED}No hay metadatos para guardar{Colors.END}")
        return
    
    original_name = metadata.get('archivo', {}).get('nombre', 'desconocido')
    json_name = f"{Path(original_name).stem}_metadatos.json"
    
    try:
        with open(json_name, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False, default=str)
        print(f"{Colors.GREEN}Metadatos guardados en: {json_name}{Colors.END}")
    except Exception as e:
        print(f"{Colors.RED}Error guardando: {e}{Colors.END}")

def generate_pdf(xonimet):
    if not xonimet.metadata:
        print(f"{Colors.RED}Primero analiza un archivo (opcion 2){Colors.END}")
        return
    
    print(f"\n{Colors.CYAN}Generando reporte PDF...{Colors.END}")
    try:
        pdf_path = xonimet.generate_pdf_report()
        if pdf_path:
            print(f"{Colors.GREEN}✓ PDF generado exitosamente{Colors.END}")
            print(f"{Colors.CYAN}📄 Ruta del reporte: {os.path.abspath(pdf_path)}{Colors.END}")
    except Exception as e:
        print(f"{Colors.RED}Error generando PDF: {e}{Colors.END}")

def modo_terminal():
    xonimet = Xonimet()
    current_file = None
    
    while True:
        clear_screen()
        print_menu()
        
        if current_file:
            print(f"{Colors.GREEN}Archivo actual: {current_file}{Colors.END}")
        else:
            print(f"{Colors.YELLOW}Ningun archivo seleccionado{Colors.END}")
        
        opcion = input(f"\n{Colors.BOLD}Selecciona una opcion [0-8]:{Colors.END} ").strip()
        
        if opcion == '1' or opcion == '6':
            new_file = select_file()
            if new_file:
                current_file = new_file
                xonimet.set_file(current_file)
                print(f"{Colors.GREEN}Archivo seleccionado: {current_file}{Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '2':
            if not current_file:
                print(f"{Colors.RED}Primero selecciona un archivo (opcion 1){Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
                continue
            
            print(f"\n{Colors.CYAN}Analizando archivo...{Colors.END}")
            xonimet.extract_all()
            clear_screen()
            xonimet.print_metadata()
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '3':
            if not xonimet.metadata:
                print(f"{Colors.RED}Primero analiza un archivo (opcion 2){Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
                continue
            save_to_json(xonimet.metadata)
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '4':
            if not xonimet.metadata:
                print(f"{Colors.RED}Primero analiza un archivo (opcion 2){Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
                continue
            generate_pdf(xonimet)
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '5':
            if not current_file:
                print(f"{Colors.RED}No hay archivo seleccionado{Colors.END}")
            else:
                print(f"\n{Colors.CYAN}Informacion del archivo actual:{Colors.END}")
                print(f"  {Colors.YELLOW}•{Colors.END} Ruta: {current_file}")
                print(f"  {Colors.YELLOW}•{Colors.END} Tamaño: {xonimet._format_bytes(os.path.getsize(current_file))}")
                print(f"  {Colors.YELLOW}•{Colors.END} Extension: {Path(current_file).suffix}")
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '7':
            clear_screen()
            print_help()
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '8':
            clear_screen()
        
        elif opcion == '0':
            print(f"\n{Colors.GREEN}Gracias por usar XONIMET 2026!{Colors.END}")
            print(f"{Colors.CYAN}Desarrollado por Darian Alberto Camacho Salas (XONIDU){Colors.END}")
            break
        
        else:
            print(f"{Colors.RED}Opcion no valida. Intenta de nuevo.{Colors.END}")
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")

# ============================================================================
# Modo Gráfico (Flask)
# ============================================================================
def modo_grafico():
    try:
        from flask import Flask, request, render_template, send_file, flash as flash_message, redirect, url_for
        import io
    except ImportError:
        print(f"{Colors.RED}Error: Flask no instalado. Ejecuta: pip install flask{Colors.END}")
        sys.exit(1)
    
    app = Flask(__name__)
    app.secret_key = "XONIMET-Darian_Alberto_Camacho_Salas"
    
    # Ruta principal
    @app.route("/", methods=["GET", "POST"])
    def index():
        if request.method == "POST":
            if 'file' not in request.files:
                flash_message('No se seleccionó ningún archivo')
                return redirect(url_for('index'))
            
            file = request.files['file']
            if file.filename == '':
                flash_message('No se seleccionó ningún archivo')
                return redirect(url_for('index'))
            
            try:
                file_content = file.read()
                filename = file.filename
                
                xonimet = Xonimet()
                xonimet.set_file_content(file_content, filename)
                metadata = xonimet.extract_all()
                
                # Formatear metadatos para mostrar
                specific = []
                spec = metadata.get('metadatos_especificos', {})
                if spec and 'error' not in spec:
                    for key, value in spec.items():
                        if isinstance(value, dict):
                            for subkey, subvalue in value.items():
                                if subvalue:
                                    specific.append({
                                        'category': key.replace('_', ' ').title(),
                                        'key': subkey.replace('_', ' ').title(),
                                        'value': str(subvalue)
                                    })
                        else:
                            if value:
                                specific.append({
                                    'category': 'Especifico',
                                    'key': key.replace('_', ' ').title(),
                                    'value': str(value)
                                })
                
                return render_template('index.html',
                                     metadata=metadata,
                                     file_info=metadata.get('archivo', {}),
                                     specific_metadata=specific,
                                     has_metadata=True,
                                     filename=filename)
            
            except Exception as e:
                flash_message(f'Error al procesar el archivo: {str(e)}')
                return redirect(url_for('index'))
        
        return render_template('index.html', has_metadata=False)
    
    # Ruta exportar JSON
    @app.route("/export_json", methods=["POST"])
    def export_json():
        try:
            metadata_json = request.form.get('metadata')
            if not metadata_json:
                flash_message('No hay metadatos para exportar')
                return redirect(url_for('index'))
            
            metadata = json.loads(metadata_json)
            json_data = json.dumps(metadata, indent=2, ensure_ascii=False, default=str)
            
            return send_file(
                io.BytesIO(json_data.encode('utf-8')),
                as_attachment=True,
                download_name=f"metadatos_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                mimetype='application/json'
            )
        except Exception as e:
            flash_message(f'Error al exportar JSON: {str(e)}')
            return redirect(url_for('index'))
    
    # Ruta exportar PDF
    @app.route("/export_pdf", methods=["POST"])
    def export_pdf():
        try:
            metadata_json = request.form.get('metadata')
            if not metadata_json:
                flash_message('No hay metadatos para exportar')
                return redirect(url_for('index'))
            
            metadata = json.loads(metadata_json)
            
            # Usar Xonimet para generar PDF
            xonimet = Xonimet()
            xonimet.metadata = metadata
            
            pdf_path = xonimet.generate_pdf_report()
            if pdf_path:
                return send_file(
                    pdf_path,
                    as_attachment=True,
                    download_name=f"reporte_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                    mimetype='application/pdf'
                )
            else:
                flash_message('Error al generar PDF')
                return redirect(url_for('index'))
        except Exception as e:
            flash_message(f'Error al exportar PDF: {str(e)}')
            return redirect(url_for('index'))
    
    # Mostrar info en terminal
    import socket
    server_ip = "127.0.0.1"
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.settimeout(0)
        try:
            s.connect(('10.254.254.254', 1))
            server_ip = s.getsockname()[0]
        except:
            pass
        finally:
            s.close()
    except:
        pass
    
    print(f"\n{Colors.GREEN}{Colors.BOLD}Servidor web iniciado{Colors.END}")
    print(f"   {Colors.CYAN}Local:{Colors.END}   http://127.0.0.1:5000/")
    print(f"   {Colors.CYAN}Red:{Colors.END}     http://{server_ip}:5000/")
    print(f"\n{Colors.YELLOW}Presiona Ctrl+C para detener el servidor{Colors.END}")
    print("-" * 50)
    
    app.run(host='0.0.0.0', port=5000, debug=False)

# ============================================================================
# Función principal
# ============================================================================
def main():
    # Limpiar pantalla
    os.system('cls' if os.name == 'nt' else 'clear')
    
    # Verificar argumentos
    if len(sys.argv) > 1:
        if sys.argv[1] in ['-h', '--help', '/?']:
            mostrar_ayuda()
            return
        
        if sys.argv[1] == '1':
            modo = 1
        elif sys.argv[1] == '2':
            modo = 2
        else:
            print(f"{Colors.RED}Error: Opcion invalida{Colors.END}")
            mostrar_ayuda()
            return
    else:
        print_banner()
        print(f"\n{Colors.BOLD}Selecciona el modo de ejecucion:{Colors.END}")
        print(f"  {Colors.CYAN}[1]{Colors.END} Modo Terminal (CLI) - Bajo consumo de recursos")
        print(f"  {Colors.CYAN}[2]{Colors.END} Modo Grafico (Web)  - Interfaz web amigable")
        print()
        
        while True:
            opcion = input(f"{Colors.GREEN}→{Colors.END} ").strip()
            if opcion in ['1', '2']:
                modo = int(opcion)
                break
            else:
                print(f"{Colors.RED}Opcion no valida. Elige 1 o 2.{Colors.END}")
    
    # Verificar dependencias
    modo_grafico_bool = (modo == 2)
    if not verificar_dependencias(modo_grafico_bool):
        print(f"\n{Colors.RED}Faltan dependencias. Instalalas y vuelve a intentarlo.{Colors.END}")
        sys.exit(1)
    
    # Ejecutar modo seleccionado
    if modo == 1:
        modo_terminal()
    else:
        try:
            modo_grafico()
        except KeyboardInterrupt:
            print(f"\n{Colors.YELLOW}Servidor detenido por el usuario{Colors.END}")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print(f"\n{Colors.YELLOW}Saliendo...{Colors.END}")
    except Exception as e:
        print(f"{Colors.RED}Error: {e}{Colors.END}")
